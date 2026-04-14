import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- Design System Constants ---
COLOR_NAVY = RGBColor(16, 42, 67)      # #102A43 (Professional Navy)
COLOR_TEAL = RGBColor(36, 59, 83)      # #243B53 (Subtle Accent)
COLOR_TEXT = RGBColor(51, 51, 51)      # #333333 (Charcoal)
COLOR_WHITE = RGBColor(255, 255, 255)
BG_IMAGE = "/gpfs-calypso/home/globc/page/.gemini/antigravity/brain/4ad20f89-f673-48d7-9ed4-28a831015efa/presentation_background_climate_professional_1776160649989.png"

def apply_premium_styling(prs, slide, title_text, is_title_slide=False):
    """Applies consistent branding and layout to a slide."""
    
    if is_title_slide:
        # Title Slide - Large Background Image
        if os.path.exists(BG_IMAGE):
            slide.shapes.add_picture(BG_IMAGE, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # Centered White Title
        title_shape = slide.shapes.title
        title_shape.left = Inches(1)
        title_shape.top = Inches(2.5)
        title_shape.width = prs.slide_width - Inches(2)
        
        tf = title_shape.text_frame
        tf.text = title_text
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
    else:
        # Content Slide - Branded Header
        header_height = Inches(0.12)
        header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, header_height)
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = COLOR_NAVY
        header_bar.line.fill.background()

        # Title Styling
        title_shape = slide.shapes.title
        title_shape.left = Inches(0.5)
        title_shape.top = Inches(0.4)
        title_shape.width = prs.slide_width - Inches(1)
        tf = title_shape.text_frame
        tf.text = title_text
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLOR_NAVY
        
        # Footer
        footer_y = prs.slide_height - Inches(0.5)
        footer_text = "EGU26 Short Course | Using Deep Learning for Climate Downscaling"
        footer = slide.shapes.add_textbox(Inches(0.5), footer_y, prs.slide_width - Inches(2), Inches(0.4))
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.text = footer_text
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEAL

        # Slide Number
        slide_num = slide.shapes.add_textbox(prs.slide_width - Inches(1), footer_y, Inches(0.5), Inches(0.4))
        tf = slide_num.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        # Placeholder for slide number logic if needed

def create_pptx(md_file, output_file):
    prs = Presentation()
    
    # 16:9 Aspect Ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    with open(md_file, 'r') as f:
        content = f.read()

    slides_raw = re.split(r'\n---', content)
    
    processed_count = 0
    for slide_text in slides_raw:
        slide_text = slide_text.strip()
        if not slide_text or slide_text.startswith('# EGU'):
            continue
            
        lines = slide_text.split('\n')
        title = "Untitled Slide"
        bullets = []
        speaker_notes = ""
        images = []
        
        in_notes = False
        for line in lines:
            line = line.strip()
            if not line: continue
            
            if re.match(r'^###?\s*Slide', line):
                title = line.split(':', 1)[1].strip() if ':' in line else line
                title = re.sub(r'\(\d+\s*min(s)?\)', '', title).strip()
                continue
                
            if line.startswith('**Speaker Notes:**'):
                in_notes = True
                speaker_notes += line.replace('**Speaker Notes:**', '').strip() + ' '
                continue
            if in_notes:
                if line.endswith('---') or line.startswith('**Keywords:**'): 
                    in_notes = False
                    continue
                speaker_notes += line + ' '
                continue
                
            img_match = re.search(r'!\[.*?\]\((.*?)\)', line)
            if img_match:
                images.append(img_match.group(1))
                continue
                
            if line.startswith('*') or line.startswith('-'):
                bullets.append(line[1:].strip())
        
        # Create Slide
        is_title = (processed_count == 0)
        layout = prs.slide_layouts[0] if is_title else prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        
        apply_premium_styling(prs, slide, title, is_title_slide=is_title)
        
        if not is_title:
            # Set Bullets
            body_shape = slide.placeholders[1]
            body_shape.left = Inches(0.8)
            body_shape.top = Inches(1.5)
            body_shape.width = prs.slide_width - Inches(5) if images else prs.slide_width - Inches(1.6)
            
            tf = body_shape.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            tf.text = "" 
            
            font_size = Pt(22)
            if len(bullets) > 6: font_size = Pt(18)
            
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = font_size
                p.font.color.rgb = COLOR_TEXT
                p.space_before = Pt(8)
            
            # Add Images
            if images:
                for i, img_path in enumerate(images):
                    if os.path.exists(img_path):
                        top = Inches(1.8 + (i * 2.6))
                        left = prs.slide_width - Inches(4)
                        slide.shapes.add_picture(img_path, left, top, height=Inches(2.5))

        # Notes
        slide.notes_slide.notes_text_frame.text = speaker_notes.strip()
        processed_count += 1
        print(f"Branded Slide {processed_count}: {title}")

    prs.save(output_file)
    print(f"\nSUCCESS: Premium presentation saved to {output_file}")

if __name__ == "__main__":
    md_path = "/scratch/globc/page/idownscale_garcia_clean/EGU_Presentation_Slides.md"
    out_path = "/scratch/globc/page/idownscale_garcia_clean/EGU26_Downscaling_ML.pptx"
    create_pptx(md_path, out_path)
